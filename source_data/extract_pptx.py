from pptx import Presentation
import sys, json

prs = Presentation('/app/source_data/aspen.pptx')
print(f"Total slides: {len(prs.slides)}")
print("="*80)

for idx, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        texts.append(run.text.strip())
        if shape.has_table:
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip())
                texts.append(" | ".join(row_text))
    if texts:
        print(f"\n--- SLIDE {idx} ---")
        for t in texts:
            print(t)
